from pathlib import Path

import unicodedata

import pandas as pd
import re
from typing import List
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from unstructured.partition.pptx import partition_pptx
from llama_index.core import SimpleDirectoryReader


class PptxLoader:
    def _load_textframed_shapes(self, shapes):
        for shape in shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_shape = shape
                    for shape in self._load_textframed_shapes(group_shape.shapes):
                        yield shape
                    continue
            except:
                continue

            if shape.has_text_frame:
                yield shape

    def _load_text(self, shapes, title_text) -> List[str]:
        all_text = []
        for shape in self._load_textframed_shapes(shapes):
            text = re.sub("\s{2,}", " ", shape.text).strip()
            if text and text != title_text:
                all_text.append(text)
        return all_text

    def _load_single_table(self, table) -> List[List[str]]:
        n_row = len(table.rows)
        n_col = len(table.columns)

        arrays = [["" for _ in range(n_row)] for _ in range(n_col)]

        for i in range(n_row):
            for j in range(n_col):
                cell = table.cell(i, j).text_frame.paragraphs
                cell_text = "".join([run.text.strip() for para in cell for run in para.runs])
                cell_text = re.sub("\n", "<br />", cell_text).strip()
                arrays[j][i] = cell_text

        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                arrays[j][i] = cell.text

        return arrays

    def _load_table(self, shapes) -> List[pd.DataFrame]:
        tables = []
        for shape in shapes:
            if shape.has_table:
                arrays = self._load_single_table(shape.table)
                tables.append(pd.DataFrame({a[0]: a[1:] for a in arrays}))
        return tables

    def _extract_content(self, presentation):
        """Extract title of slide,  text body and table elements from presentation

        Args:
            presentation (pptx.Presentation):

        Returns:
            list[str]: a list contains title of slide
            list[str]: a list contains text body
            list[str]: a list contains table content
        """
        all_title = []
        all_text = []
        tables = []
        for slide_number, slide in enumerate(presentation.slides):
            # -- extract slide title
            title_text = ""
            if (title_shape := slide.shapes.title) is not None:
                title_text = title_shape.text.strip()
            all_title.append(title_text)

            try:
                sorted_shapes = sorted(slide.shapes, key=lambda shape: (shape.top, shape.left))
            except Exception:
                sorted_shapes = slide.shapes

            # -- extract text
            slide_text = self._load_text(sorted_shapes, title_text)
            slide_text = "\n".join(
                [unicodedata.normalize("NFKC", text) for text in slide_text]
            )
            all_text.append(slide_text)

            # -- extract table
            slide_tables = self._load_table(sorted_shapes)
            tables.append(slide_tables)

        return all_title, all_text, tables

    def load_data(self, file_path: str) -> List[str]:
        """Load data using Pptx reader, considering only text and tables

        Args:
            file_path (str):

        Returns:
            List[Document]: list of documents extracted from file
        """
        file_path = Path(file_path).resolve()

        doc = Presentation(str(file_path))
        all_title, all_text, tables = self._extract_content(doc)

        # create output Document with metadata from table
        documents = [
            table.to_csv(index=False).strip()
            for slide_tables in tables
            for table in slide_tables
        ]

        # create Document from non-table text
        documents.extend(
            [non_table_text.strip()
             for non_table_text in all_text]
        )

        return documents


class UnstructuredPptxReader:
    """Load data using unstructured library
    Ref: https://docs.unstructured.io/open-source/core-functionality/partitioning#partition-pptx
    """

    def load_data(self, file_path: str) -> List[str]:
        elements = partition_pptx(filename=file_path)
        all_text = "\n".join([ele.text.strip() for ele in elements])
        return [all_text]


class LlamaIndexPptxReader:
    """Load data using llama-index library
    Ref: https://docs.llamaindex.ai/en/stable/module_guides/loading/simpledirectoryreader/
    """

    def load_data(self, file_path: str) -> List[str]:
        reader = SimpleDirectoryReader(input_files=[file_path])
        document = reader.load_data()[0]
        return [document.text.strip()]
