"""Power point interface."""

from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.util import Inches, Length
from typing import Optional


def add_background_image(
    slide: Slide,
    image_file: str,
    left: Optional[Length] = 0,
    top: Optional[Length] = 0,
    width: Optional[Length] = None,
    height: Optional[Length] = None,
    move_to_background: Optional[bool] = True,
) -> Picture:
    """Add background image to a slide.

    Parmeters:
        move_to_background: If `True`, will move image furtherst back in the stack of elements.
    """
    background_img = slide.shapes.add_picture(image_file=image_file, left=left, top=top, width=width, height=height)
    if move_to_background:
        slide.shapes._spTree.remove(background_img._element)
        slide.shapes._spTree.insert(2, background_img._element)
    return background_img


class PowerPoint:
    """PowerPoint class with pptx.Presentation as `prs` attribute.

    Attributes:
        `prs`: Contains the pptx.Presentation object.

    Methods:
        `add_slide`: Add a slide from default tempalte using template index and having an optional title.
    """

    widescreen_width = Inches(13.33)
    widescreen_height = Inches(7.5)

    def __init__(self) -> None:
        self.prs = Presentation()

    def add_slide(self, layout_index: int, title_text: Optional[str] = None) -> Slide:
        """Add a PowerPoint slide with an optional title text.

        Returns:
            slide object
        """
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        if title_text:
            title = slide.shapes.title
            title.text = title_text
        return slide


def default_templates(widescreen_dimensions: Optional[bool] = True) -> None:
    """Create a power point slide presentation using default templates."""
    ppt = PowerPoint()

    if widescreen_dimensions:
        ppt.prs.slide_width = PowerPoint.widescreen_width
        ppt.prs.slide_height = PowerPoint.widescreen_height

    [ppt.add_slide(layout_index) for layout_index in range(11)]

    ppt.prs.save("default_templates.pptx")
