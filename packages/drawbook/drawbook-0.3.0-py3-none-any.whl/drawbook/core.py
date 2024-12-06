"""
Core functionality for the drawbook library.
"""

from pathlib import Path
from typing import List, Literal
import tempfile
import io
import requests
import warnings
from tqdm import tqdm
from PIL import Image
import huggingface_hub
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from huggingface_hub import InferenceClient
from PIL import ImageDraw, ImageFont
import gradio as gr
import textwrap
import json


class Book:
    """A class representing a children's book that can be exported to PowerPoint."""

    def __init__(
        self,
        title: str = "Untitled Book",
        pages: List[str] = None,
        title_illustration: str | Literal[False] | None = None,
        illustrations: List[str | None | Literal[False]] = None,
        lora: str = "SebastianBodza/Flux_Aquarell_Watercolor_v2",
        author: str | None = None,
        illustration_prompts: List[str | None] = None,
        title_illustration_prompt: str | None = None,
    ):
        """
        Initialize a new Book.

        Args:
            title: The book's title
            pages: List of strings containing text for each page
            illustrations: List of illustration paths or placeholders
                         (str for path, None for pending, False for no illustration)
            lora: The LoRA model on Hugging Face to use for illustrations
            author: The book's author name
            illustration_prompts: Optional list of custom prompts for page illustrations
            title_illustration_prompt: Optional custom prompt for title illustration
        """
        self.title = title
        self.pages = pages or []
        self.illustrations = illustrations or []
        self.lora = lora
        self.title_illustration = title_illustration
        self.author = author
        self.illustration_prompts = illustration_prompts or []
        self.title_illustration_prompt = title_illustration_prompt
        self.client = InferenceClient()
        self.page_previews = []

        # Ensure illustrations list matches pages length
        while len(self.illustrations) < len(self.pages):
            self.illustrations.append(None)

        # Ensure illustration_prompts list matches pages length
        while len(self.illustration_prompts) < len(self.pages):
            self.illustration_prompts.append(None)

    def _get_illustration_prompt(self, text: str) -> str:
        """Get an illustration prompt from the text using Qwen."""
        system_prompt = """You are a helpful assistant that converts children's book text into illustration prompts. 
        Extract a key object along with its description that could be used to illustrate the page. 
        Replace any proper names with more generic versions.
        
        For example:
        If the text is: "Mustafa loves his silver cybertruck. One day, his cybertruck starts to glow, grow, and zoom up into the sky"
        You should return: "A silver cybertruck zooming into the sky"
        
        If the text is: "Up, up, up goes Mustafa in his special cybertruck. He waves bye-bye to his house as it gets tiny down below"
        You should return: "A boy in the sky waving bye"
        """

        user_prompt = f"""This is the text of a page in a children's book. From this text, extract a key object along with its description that could be used to illustrate this page. Replace any proper names with more generic versions.

Text: {text}

Return ONLY the illustration description, nothing else."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        try:
            stream = self.client.chat.completions.create(
                model="Qwen/Qwen2.5-72B-Instruct",
                messages=messages,
                max_tokens=500,
                stream=True,
            )

            response = ""
            for chunk in stream:
                if chunk.choices[0].delta.content is not None:
                    response += chunk.choices[0].delta.content

            return response.strip()
        except Exception:
            gr.warning("Could not access Hugging Face Inference API, make sure that you are logged in locally to your Hugging Face account")
            return text            

    def _get_prompt(self, illustration_prompt: str) -> str:
        if self.lora == "SebastianBodza/Flux_Aquarell_Watercolor_v2":
            return f"A AQUACOLTOK watercolor painting with a white background of: {illustration_prompt}"
        else:
            warnings.warn(
                f"The LoRA model '{self.lora}' is not officially supported. "
                "Results may not be as expected.",
                UserWarning,
            )
            return f"An illustration of: {illustration_prompt}"

    def export(self, filename: str | Path | None = None) -> None:
        """
        Export the book to a PowerPoint file.

        Args:
            filename: Optional path where to save the file. If None, creates in temp directory.
        """
        if filename is None:
            # Create temp file with .pptx extension
            temp_file = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
            output_path = Path(temp_file.name)
            temp_file.close()
        else:
            # Convert to Path object and resolve to absolute path
            output_path = Path(filename).resolve()
            # Ensure parent directories exist
            output_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        # Add diagonal striped border at the top
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor(128, 0, 0)  # Maroon

        # Add title illustration first if available
        if isinstance(self.title_illustration, str):
            try:
                slide.shapes.add_picture(
                    self.title_illustration,
                    Inches(2.5),
                    Inches(1.5),
                    Inches(5),
                    Inches(5),
                )
            except Exception as e:
                print(f"Warning: Could not add title illustration: {e}")

        # Add title with adjusted positioning and z-order
        title = slide.shapes.title
        title.top = Inches(0)
        title.height = Inches(2.0)
        title.width = Inches(10)

        # Add title text
        p1 = title.text_frame.paragraphs[0]
        # Clear any existing text
        p1.clear()
        p1.font.name = "Trebuchet MS"
        p1.alignment = PP_ALIGN.CENTER

        # Define common stop words
        stop_words = {
            "a",
            "an",
            "and",
            "are",
            "as",
            "at",
            "be",
            "by",
            "for",
            "from",
            "has",
            "he",
            "in",
            "is",
            "it",
            "its",
            "of",
            "on",
            "that",
            "the",
            "to",
            "was",
            "were",
            "will",
            "with",
        }

        # Split title and add each word with appropriate size
        words = self.title.split()
        for i, word in enumerate(words):
            run = p1.add_run()
            run.text = word + (" " if i < len(words) - 1 else "")
            run.font.name = "Trebuchet MS"
            if word.lower() in stop_words:
                run.font.size = Inches(0.42)  # Smaller size for stop words
            else:
                run.font.size = Inches(0.5)  # Regular size for other words

        # Add author with adjusted positioning
        if self.author is not None:
            author_box = slide.shapes.add_textbox(
                Inches(0),
                Inches(6.5),  # Moved up from bottom
                Inches(10),
                Inches(0.5),
            )
            author_frame = author_box.text_frame
            author_frame.text = f"Written by {self.author}"
            author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            author_frame.paragraphs[0].font.name = "Trebuchet MS"
            author_frame.paragraphs[0].font.size = Inches(0.25)

        # Add content slides
        content_slide_layout = prs.slide_layouts[5]  # Blank layout

        for page_num, (text, illustration) in enumerate(
            zip(self.pages, self.illustrations)
        ):
            slide = prs.slides.add_slide(content_slide_layout)

            if isinstance(illustration, str):
                try:
                    slide.shapes.add_picture(
                        illustration, Inches(2.5), Inches(1.4), Inches(5), Inches(5)
                    )
                except Exception as e:
                    print(
                        f"Warning: Could not add illustration on page {page_num + 1}: {e}"
                    )

            # Split text into sentences and join with newlines
            sentences = text.replace(". ", ".\n").split("\n")

            # Special formatting for first page
            if page_num == 0 and text:
                # Split first character from first sentence
                first_char = sentences[0][0]
                first_sentence_rest = sentences[0][1:]

                p = slide.shapes.title.text_frame.paragraphs[0]
                p.line_spacing = 1.5  # Add line spacing
                run = p.add_run()
                run.text = first_char
                run.font.size = Inches(0.3)
                run.font.name = "Trebuchet MS"

                run = p.add_run()
                run.text = first_sentence_rest
                run.font.size = Inches(0.25)
                run.font.name = "Trebuchet MS"

                # Add remaining sentences as new paragraphs
                for sentence in sentences[1:]:
                    p = slide.shapes.title.text_frame.add_paragraph()
                    p.line_spacing = 1.5  # Add line spacing
                    p.text = sentence
                    p.font.name = "Trebuchet MS"
                    p.font.size = Inches(0.25)
                    p.alignment = PP_ALIGN.CENTER
            else:
                # Add each sentence as a separate paragraph
                first_paragraph = True
                for sentence in sentences:
                    if first_paragraph:
                        p = slide.shapes.title.text_frame.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = slide.shapes.title.text_frame.add_paragraph()
                    p.line_spacing = 1.5  # Add line spacing
                    p.text = sentence
                    p.font.name = "Trebuchet MS"
                    p.font.size = Inches(0.25)
                    p.alignment = PP_ALIGN.CENTER

            # Add page number at bottom center
            page_number = page_num + 1  # Add 1 since page_num is 0-based
            page_num_box = slide.shapes.add_textbox(
                Inches(0),
                Inches(6.5),  # Y position near bottom of slide
                Inches(10),
                Inches(0.5),  # Full width of slide for centering
            )
            page_num_frame = page_num_box.text_frame
            page_num_frame.text = str(page_number)
            page_num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            page_num_frame.paragraphs[0].font.name = "Trebuchet MS"
            page_num_frame.paragraphs[0].font.size = Inches(
                0.15
            )  # Slightly smaller than main text

        # Save the presentation
        prs.save(str(output_path))
        print(f"Book exported to: {output_path.absolute()}")
        return output_path.absolute()

    def __len__(self) -> int:
        """Return the number of pages in the book."""
        return len(self.pages)

    def illustrate(
        self, save_dir: str | Path | None = None, page_num: int | None = None
    ) -> str | None:
        """
        Generate illustrations using the Hugging Face Inference API.

        Args:
            save_dir: Optional directory to save the generated images.
                     If None, creates a temporary directory.
            page_num: Optional specific page to illustrate (0 for title page, 1+ for content pages).
                     If None, illustrates all pages.

        Returns:
            Status message if page_num is specified, None otherwise.
        """
        token = huggingface_hub.get_token()
        if not token:
            msg = "No Hugging Face token found. Please login using `huggingface-cli login`"
            if page_num is not None:
                return f"Error: {msg}"
            warnings.warn(msg)

        API_URL = f"https://api-inference.huggingface.co/models/{self.lora}"
        headers = {"Authorization": f"Bearer {token}"}

        # Create save directory if provided
        if save_dir:
            save_dir = Path(save_dir)
            save_dir.mkdir(parents=True, exist_ok=True)
        else:
            save_dir = Path(tempfile.mkdtemp())

        # Create list of tasks
        if page_num is not None:
            if page_num == 0:
                tasks = [("title", self.title, self.title_illustration)]
            else:
                page_idx = page_num - 1
                tasks = [
                    (
                        f"page_{page_num}",
                        self.pages[page_idx],
                        self.illustrations[page_idx],
                    )
                ]
        else:
            print("Generating illustrations... This could take a few minutes.")
            tasks = []
            if self.title_illustration is None:
                tasks.append(("title", self.title, None))
            tasks.extend(
                (f"page_{i+1}", text, current_illust)
                for i, (text, current_illust) in enumerate(
                    zip(self.pages, self.illustrations)
                )
            )

        for task_name, text, current_illust in tqdm(
            tasks, desc="Generating illustrations", disable=page_num is not None
        ):
            # Skip if illustration already exists or is explicitly disabled
            if isinstance(current_illust, str) or current_illust is False:
                continue

            try:
                if page_num is None:
                    print(f"\n=== Processing {task_name} ===")
                    print(f"Original text: {text}")

                if task_name == "title":
                    if not self.title_illustration_prompt:
                        self.title_illustration_prompt = self._get_illustration_prompt(
                            text
                        )
                    if page_num is None:
                        print(
                            f"Title illustration prompt: {self.title_illustration_prompt}"
                        )
                    prompt = self._get_prompt(self.title_illustration_prompt)
                    if page_num is None:
                        print(f"Final title image prompt: {prompt}")
                else:
                    page_idx = int(task_name.split("_")[1]) - 1
                    if not self.illustration_prompts[page_idx]:
                        self.illustration_prompts[page_idx] = (
                            self._get_illustration_prompt(text)
                        )
                    if page_num is None:
                        print(
                            f"Illustration prompt: {self.illustration_prompts[page_idx]}"
                        )
                    prompt = self._get_prompt(self.illustration_prompts[page_idx])
                    if page_num is None:
                        print(f"Final image prompt: {prompt}")

                response = requests.post(
                    API_URL, headers=headers, json={"inputs": prompt}
                )

                if response.status_code != 200:
                    msg = f"Failed to generate illustration for {task_name}: {response.text}"
                    if page_num is not None:
                        return f"Error: {msg}"
                    print(f"Warning: {msg}")
                    continue

                # Save the image
                image = Image.open(io.BytesIO(response.content))
                image_path = save_dir / f"{task_name}.png"
                image.save(image_path)
                if page_num is None:
                    print(f"Image saved to: {image_path}")

                # Update the appropriate illustration reference
                if task_name == "title":
                    self.title_illustration = str(image_path)
                else:
                    page_idx = int(task_name.split("_")[1]) - 1
                    self.illustrations[page_idx] = str(image_path)

            except Exception as e:
                msg = f"Error generating illustration for {task_name}: {e}"
                if page_num is not None:
                    return f"Error: {msg}"
                print(f"Warning: {msg}")
                continue

        if page_num is None:
            print(f"\nAll illustrations saved to: {save_dir}")
        else:
            return "Illustration generated successfully!"

    def create_preview(self, page_num: int | None = None):
        """
        Create visual previews of book pages.

        Args:
            page_num: Optional specific page to preview (0 for title page, 1+ for content pages).
                     If None, creates previews for all pages.
        """
        # Constants for page layout (matching PowerPoint dimensions and positioning)
        PAGE_WIDTH = 1920
        PAGE_HEIGHT = 1080
        ILLUSTRATION_WIDTH = 840
        ILLUSTRATION_HEIGHT = 840
        ILLUSTRATION_X = (PAGE_WIDTH - ILLUSTRATION_WIDTH) // 2
        ILLUSTRATION_Y = 120

        # Try to load Trebuchet MS font, fall back to Arial if not available
        try:
            title_font = ImageFont.truetype("Trebuchet MS", 96)
            body_font = ImageFont.truetype("Trebuchet MS", 48)
            page_num_font = ImageFont.truetype("Trebuchet MS", 29)
            author_font = ImageFont.truetype("Trebuchet MS", 48)
        except OSError:
            title_font = ImageFont.truetype("Arial", 96)
            body_font = ImageFont.truetype("Arial", 48)
            page_num_font = ImageFont.truetype("Arial", 29)
            author_font = ImageFont.truetype("Arial", 48)

        # Determine which pages to process
        if page_num is not None:
            if page_num == 0:
                pages_to_process = [(0, "title", None, None)]
            else:
                page_idx = page_num - 1
                pages_to_process = [
                    (
                        page_num,
                        "content",
                        self.pages[page_idx],
                        self.illustrations[page_idx],
                    )
                ]
        else:
            # Process all pages
            pages_to_process = [(0, "title", None, None)]  # Title page
            pages_to_process.extend(
                (i + 1, "content", text, illust)
                for i, (text, illust) in enumerate(zip(self.pages, self.illustrations))
            )

        # Process each page
        for page_num, page_type, text, illustration in pages_to_process:
            if page_type == "title":
                # Create title page
                page = Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), "white")
                draw = ImageDraw.Draw(page)

                # Add maroon border on left
                draw.rectangle([(0, 0), (38, PAGE_HEIGHT)], fill=(128, 0, 0))

                # Add title illustration if available
                if isinstance(self.title_illustration, str):
                    try:
                        illust = Image.open(self.title_illustration)
                        illust = illust.resize(
                            (ILLUSTRATION_WIDTH, ILLUSTRATION_HEIGHT)
                        )
                        page.paste(illust, (ILLUSTRATION_X, ILLUSTRATION_Y))
                    except Exception as e:
                        print(f"Warning: Could not add title illustration: {e}")

                # Add title text
                title_y = 40
                bbox = draw.textbbox((0, 0), self.title, font=title_font)
                title_width = bbox[2] - bbox[0]
                draw.text(
                    ((PAGE_WIDTH - title_width) // 2, title_y),
                    self.title,
                    font=title_font,
                    fill="black",
                )

                # Add author if available
                if self.author:
                    author_text = f"Written by {self.author}"
                    bbox = draw.textbbox((0, 0), author_text, font=author_font)
                    author_width = bbox[2] - bbox[0]
                    draw.text(
                        ((PAGE_WIDTH - author_width) // 2, PAGE_HEIGHT - 100),
                        author_text,
                        font=author_font,
                        fill="black",
                    )

            else:
                # Create content page
                page = Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), "white")
                draw = ImageDraw.Draw(page)

                # Add illustration if available
                if isinstance(illustration, str):
                    try:
                        illust = Image.open(illustration)
                        illust = illust.resize(
                            (ILLUSTRATION_WIDTH, ILLUSTRATION_HEIGHT)
                        )
                        page.paste(illust, (ILLUSTRATION_X, ILLUSTRATION_Y))
                    except Exception as e:
                        print(
                            f"Warning: Could not add illustration on page {page_num}: {e}"
                        )

                # Add text
                text_y = 50
                wrapped_text = textwrap.fill(text, width=50)
                bbox = draw.textbbox((0, 0), wrapped_text, font=body_font)
                text_width = bbox[2] - bbox[0]
                draw.text(
                    ((PAGE_WIDTH - text_width) // 2, text_y),
                    wrapped_text,
                    font=body_font,
                    fill="black",
                    align="center",
                )

                # Add page number
                page_num_text = str(page_num)
                bbox = draw.textbbox((0, 0), page_num_text, font=page_num_font)
                page_num_width = bbox[2] - bbox[0]
                draw.text(
                    ((PAGE_WIDTH - page_num_width) // 2, PAGE_HEIGHT - 100),
                    page_num_text,
                    font=page_num_font,
                    fill="black",
                )

            # Update the preview pages
            if page_num is None:
                self.page_previews.append(page)
            else:
                # Ensure list is long enough
                while len(self.page_previews) <= page_num:
                    self.page_previews.append(None)
                self.page_previews[page_num] = page

        return self.page_previews if page_num is None else self.page_previews[page_num]

    def preview(self) -> None:
        """
        Create a visual preview of the book pages and display them in a Gradio interface.
        """
        print("Creating preview...")
        self.create_preview()

        # Launch Gradio interface
        with gr.Blocks(theme="citrus") as preview_interface:
            selected_page = gr.State(0)

            def select_page(selected: gr.SelectData):
                index = selected.index
                if index == 0:
                    return self.title, self.title_illustration_prompt, index
                else:
                    return (
                        self.pages[index - 1],
                        self.illustration_prompts[index - 1],
                        index,
                    )

            def export_book():
                output_path = self.export()
                return gr.DownloadButton(value=output_path, interactive=True)

            gr.Markdown(f"<center><h1>{self.title}</h1></center>")
            with gr.Row():
                with gr.Column():
                    page = gr.Textbox(self.title, label="Page text", lines=3, interactive=False)
                    prompt = gr.Textbox(
                        self.title_illustration_prompt,
                        label="Illustration prompt",
                        lines=3,
                    )
                    with gr.Row():
                        prompt_button = gr.Button(
                            "Generate Prompt", variant="secondary"
                        )
                        image_button = gr.Button("Generate Image", variant="primary")
                with gr.Column():
                    gallery = gr.Gallery(
                        value=self.page_previews,
                        columns=2,
                        rows=2,
                        height=600,
                        show_label=False,
                        preview=True,
                    )
                    with gr.Row():
                        export_button = gr.Button("Export", variant="secondary")
                        download_button = gr.DownloadButton(
                            label="Download", variant="primary", interactive=False
                        )

            def generate_prompt_page(selected_page: int, page_text: str):
                yield {prompt_button: gr.Button("Generating...", interactive=False)}
                illustration_prompt = self._get_illustration_prompt(page_text)
                if selected_page == 0:
                    self.title = page_text
                    self.title_illustration_prompt = illustration_prompt
                else:
                    self.pages[selected_page - 1] = page_text
                    self.illustration_prompts[selected_page - 1] = illustration_prompt
                yield {prompt_button: gr.Button("Generate Prompt", interactive=True), prompt: illustration_prompt}

            def generate_illustration_page(selected_page: int, page_text: str, illustration_prompt: str):
                yield {image_button: gr.Button("Generating...", interactive=False)}
                if not illustration_prompt:
                    illustration_prompt = self._get_illustration_prompt(page_text)
                    yield {prompt: illustration_prompt}
                self.illustrate(page_num=selected_page)
                self.create_preview(page_num=selected_page)
                print("self.page_previews", self.page_previews)
                yield {gallery: self.page_previews, image_button: gr.Button("Generate Image", interactive=True)}

            gallery.select(
                select_page,
                outputs=[page, prompt, selected_page],
                show_progress="hidden",
            )
            prompt_button.click(
                fn=generate_prompt_page,
                inputs=[selected_page, page],
                outputs=[prompt_button, prompt],
                show_progress="minimal",
            )
            image_button.click(
                fn=generate_illustration_page,
                inputs=[selected_page, page, prompt],
                outputs=[image_button, gallery, prompt],
                show_progress="minimal",
            )
            export_button.click(
                fn=export_book,
                inputs=[],
                outputs=[download_button],
                show_progress="minimal",
            )

        preview_interface.launch()

    def save(self, filepath: str | Path = "book.json") -> None:
        """
        Save the book data to a JSON file.

        Args:
            filepath: Path where to save the JSON file. Defaults to "book.json"
        """
        filepath = Path(filepath)

        # Create dictionary of book data
        book_data = {
            "title": self.title,
            "pages": self.pages,
            "title_illustration": self.title_illustration,
            "illustrations": self.illustrations,
            "lora": self.lora,
            "author": self.author,
            "illustration_prompts": self.illustration_prompts,
            "title_illustration_prompt": self.title_illustration_prompt,
        }

        # Save to JSON file
        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(book_data, f, indent=2, ensure_ascii=False)

        print(f"Book saved to: {filepath.absolute()}")

    @classmethod
    def load(cls, filepath: str | Path = "book.json") -> "Book":
        """
        Load a book from a JSON file.

        Args:
            filepath: Path to the JSON file to load. Defaults to "book.json"

        Returns:
            A new Book instance with the loaded data
        """
        filepath = Path(filepath)

        if not filepath.exists():
            raise FileNotFoundError(f"No book file found at: {filepath}")

        with open(filepath, "r", encoding="utf-8") as f:
            book_data = json.load(f)

        # Create new book instance with loaded data
        book = cls(
            title=book_data["title"],
            pages=book_data["pages"],
            title_illustration=book_data["title_illustration"],
            illustrations=book_data["illustrations"],
            lora=book_data["lora"],
            author=book_data["author"],
            illustration_prompts=book_data["illustration_prompts"],
            title_illustration_prompt=book_data["title_illustration_prompt"],
        )

        return book
