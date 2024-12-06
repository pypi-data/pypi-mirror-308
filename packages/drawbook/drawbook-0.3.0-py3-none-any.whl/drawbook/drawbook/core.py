"""
Core functionality for the drawbook library.
"""

from pathlib import Path
from typing import List, Optional, Union
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ... existing imports and class definition ...

    def _get_title_font(self) -> str:
        """Get the best available decorative font for the title."""
        # List of fonts in order of preference
        title_fonts = [
            "Palatino",  # Available on Mac
            "Baskerville",  # Available on Mac
            "Gabriola",  # Available on Windows
            "Times New Roman",  # Universal fallback
        ]
        return title_fonts[0]  # For now just return first option
        # TODO: Add font availability checking

    def _get_body_font(self) -> str:
        """Get the best available serif font for the body text."""
        # List of fonts in order of preference
        body_fonts = [
            "Palatino",  # Available on Mac
            "Georgia",  # Common on both platforms
            "Times New Roman",  # Universal fallback
        ]
        return body_fonts[0]  # For now just return first option
        # TODO: Add font availability checking

    def export(self) -> None:
        """
        Export the book to a PowerPoint file in a temporary location and print the path.
        """
        # Create temp file with .pptx extension
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        output_path = Path(temp_file.name)
        temp_file.close()

        prs = Presentation()
        
        # Add title slide with custom styling
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = self.title
        
        # Style the title with a decorative font
        title_paragraph = title.text_frame.paragraphs[0]
        title_run = title_paragraph.runs[0]
        title_run.font.name = self._get_title_font()
        title_run.font.size = Pt(72)
        title_run.font.bold = True
        
        # Add content slides
        content_slide_layout = prs.slide_layouts[5]  # Blank layout
        body_font = self._get_body_font()
        
        for page_num, (text, illustration) in enumerate(zip(self.pages, self.illustrations)):
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Add text
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                Inches(9), Inches(1)
            )
            tf = txBox.text_frame
            
            # Style the text differently for the first page
            if page_num == 0:
                # Create first paragraph with drop cap
                first_char = text[0]
                rest_of_text = text[1:]
                
                p = tf.paragraphs[0]
                drop_cap_run = p.add_run()
                drop_cap_run.text = first_char
                drop_cap_run.font.name = body_font
                drop_cap_run.font.size = Pt(60)
                drop_cap_run.font.bold = True
                
                regular_run = p.add_run()
                regular_run.text = rest_of_text
                regular_run.font.name = body_font
                regular_run.font.size = Pt(24)
            else:
                # Regular pages
                tf.text = text
                p = tf.paragraphs[0]
                p.runs[0].font.name = body_font
                p.runs[0].font.size = Pt(24)
            
            # Add illustration if available
            if isinstance(illustration, str):
                try:
                    slide.shapes.add_picture(
                        illustration,
                        Inches(1), Inches(2),
                        Inches(8), Inches(4)
                    )
                except Exception as e:
                    print(f"Warning: Could not add illustration on page {page_num + 1}: {e}")
        
        # Save the presentation
        prs.save(str(output_path))
        print(f"Book exported to: {output_path.absolute()}") 