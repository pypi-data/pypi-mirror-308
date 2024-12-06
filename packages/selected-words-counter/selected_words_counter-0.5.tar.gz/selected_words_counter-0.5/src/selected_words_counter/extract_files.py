import os
import re
import shutil
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from glob import glob

import extract_msg
import pandas as pd
import pptx
import textract
from docx import Document
from pypdf import PdfReader
from tqdm import tqdm

"""

This file is used to store code for extracting text from different file formats.

@author: Michael de Winter
"""


def process_and_save_file(
    afilepath,
    atarget_dir,
    atarget_dir_extracted,
    verbose=False,
):
    """

    Method to call process_file and correctly generate a filename if it contains subdirectories.

    @param afilepath: Path to a specific file.
    @param atarget_dir: the original directory of the file is stripped from the filename for naming purposes to get it's subdirectory.
    @oaram atarget_dir_extracted: directory in which the file will be stored.
    """

    if verbose:
        print("-------")
        print(afilepath)
    afilepath = afilepath.replace("\\", "/")

    # remove linebreaks with "-"
    text_content = re.sub(r"-\s*\n\s*", "", process_file(afilepath))
    # Normalize a file by removing \b and whitespaces
    text_content = re.sub(r"\s+", " ", text_content)

    try:
        # Generate a output name but if there are subdirectories in the directory put the subdirectories name in the filename.
        a_output_name = (
            atarget_dir_extracted
            + "/"
            + re.sub(
                r"\.(?!.*\.)",
                "_",
                afilepath.replace(atarget_dir, "").replace("/", "#"),
            )
            + ".txt"
        )
        if verbose:
            print(a_output_name)
        with open(a_output_name, "w", encoding="utf-8") as file:
            file.write(str(text_content))
    except Exception as e:
        print("Error with writing file: " + e)


def make_dir_from_filename(afilepath):
    """
    Make a directory based on the file name

    """
    a_output_directory = afilepath.rsplit(".", 1)[0]
    print("Making :" + str(a_output_directory))
    os.makedirs(a_output_directory)
    return a_output_directory


def extract_msg_attachments(atarget_dir):
    """

    Extract all .msg files in a directory.

    @param atarget_dir: directory where the .msg files are stored.
    """
    # Extract all .msg files into a directory
    for afilepath in glob(atarget_dir + "*.msg"):
        afilepath = afilepath.replace("\\", "/")

        encodings = ["utf-8", "ISO-8859-1", "windows-1252"]
        for encoding in encodings:
            try:
                msg = extract_msg.Message(afilepath, overrideEncoding=encoding)

                if len(msg.attachments) > 0:
                    a_output_directory = make_dir_from_filename(afilepath)

                    for item in range(0, len(msg.attachments)):
                        att = msg.attachments[item]
                        msg.attachments[item].save(
                            customPath=a_output_directory,
                            customFilename=att.longFilename,
                        )
                return 0
            except Exception as e:
                continue
        print(f"Failed to read {afilepath} with all tested encodings.")
        return None


def extract_zip_attachments(atarget_dir):
    """

    Extract all .zip files in a directory.

    @param atarget_dir: directory where the .zip files are stored.
    """
    # Extract a zip file into a new directory

    for afilepath in glob(atarget_dir + "*.zip"):
        try:
            afilepath = afilepath.replace("\\", "/")
            print(afilepath)
            a_output_directory = make_dir_from_filename(afilepath)

            # Extract the contents of the zip file
            with zipfile.ZipFile(afilepath, "r") as zip_ref:
                zip_ref.extractall(a_output_directory)
        except Exception as e:
            print(e)


def extracted_files_from_list_filepaths(
    afilepaths, atarget_dir, atarget_dir_extracted, verbose=False, threads=False
):
    """

    Method to extract all files in directory in a .txt format in a other directory for easy searching and snapshotting.

    @param afilepaths: a array of filepaths to be extracted to .txt files.
    @param atarget_dir: the original directory of the file is stripped from the filename for naming purposes to get it's subdirectory.
    @oaram atarget_dir_extracted: directory in which the file will be stored.
    @param verbose: Whether to print more output.
    @param threads: Whether to multiprocess with the file reading thus speeding up the process but not all systems can handle it.
    """
    if threads:
        with ThreadPoolExecutor() as executor:
            futures = {
                executor.submit(
                    process_and_save_file,
                    afilepath,
                    atarget_dir,
                    atarget_dir_extracted,
                    verbose,
                ): afilepath
                for afilepath in afilepaths
            }
            for future in tqdm(as_completed(futures), total=len(futures)):
                try:
                    future.result()
                except Exception as e:
                    print(f"Exception occurred: {e}")
    else:
        for afilepath in tqdm(afilepaths):
            process_and_save_file(
                afilepath, atarget_dir, atarget_dir_extracted, verbose
            )


def run(atarget_dir, atarget_dir_extracted, amulti_thread=False):
    """
    Method to run the process.

    @param atarget_dir: the original directory of the file is stripped from the filename for naming purposes to get it's subdirectory.
    @oaram atarget_dir_extracted: directory in which the file will be stored.
    @param amulti_thread: Whether to multiprocess with the file reading thus speeding up the process but not all systems can handle it. TODO: Find out why
    """
    # First extract all the .msg files.
    print("Extracting .msg files")
    extract_msg_attachments(atarget_dir)
    # Then extract all the .zip files.
    print("Extracting .zip files")
    extract_zip_attachments(atarget_dir)

    # Make a directory if the directory does not exist yet.
    if os.path.isdir(atarget_dir_extracted) == False:
        os.makedirs(atarget_dir_extracted)

    extracted_files_from_list_filepaths(
        [
            afilepath
            for afilepath in glob(os.path.join(atarget_dir, "**", "*"), recursive=True)
        ],
        atarget_dir,
        atarget_dir_extracted,
        threads=amulti_thread,
    )


# Different functions for opening files.
def read_pdf(file_path):
    text = ""
    try:
        # Open the PDF file
        with open(file_path, "rb") as file:
            # Create a PDF reader object
            pdf_reader = PdfReader(file)

            # Get the number of pages in the PDF
            num_pages = len(pdf_reader.pages)

            # Extract text from each page
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text += (
                    page.extract_text()
                    .encode("utf-8", errors="ignore")
                    .decode("utf-8", errors="ignore")
                )

    except Exception as e:
        print(f"Error processing {file_path}: {e}")

    return text


def read_docx(file_path):
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)


def read_msg(file_path):
    encodings = ["utf-8", "ISO-8859-1", "windows-1252"]
    for encoding in encodings:
        try:
            msg = extract_msg.Message(file_path, overrideEncoding=encoding)
            msg_text = "Sender: " + str(msg.sender) + " | \n "
            msg_text += "To: " + str(msg.to) + " | \n "
            msg_text += "CC: " + str(msg.cc) + " | \n "
            msg_text += "BCC: " + str(msg.bcc) + " | \n "
            msg_text += "Subject: " + str(msg.subject) + " | \n "
            msg_text += "Body: " + str(msg.body)
            return msg_text  # Return if successful
        except Exception as e:
            continue  # Try the next encoding if an error occurs

    print(f"Failed to read {file_path} with all tested encodings.")
    return None  # Return None if all encoding attempts fail


def read_xls(file_path):
    # Convert the pandas DataFrame to a single string
    data = pd.read_excel(file_path)
    collapsed_text = data.to_string(index=False)
    return collapsed_text


def read_doc(file_path):
    text = textract.process(file_path)
    return text


def read_pptx(file_path):
    text = ""
    presentation = pptx.Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text


def process_file(a_file_path):
    """

    Open a file based on a file extension.

    @param a_file_path: File path to a file.
    """

    file_split = a_file_path.split(".")
    file_type = a_file_path.rsplit(".", 1)[-1]
    a_content = ""

    if len(file_split) > 1:
        try:
            if "pdf" in file_type:
                # Action for PDF file type
                a_content = read_pdf(a_file_path)
            elif "xls" in file_type or "xlsx" in file_type:
                # Action for XLS file type
                a_content = read_xls(a_file_path)
            elif "docx" in file_type:
                # Action for DOCX file type
                a_content = read_docx(a_file_path)
            elif "doc" in file_type:
                a_content = read_doc(a_file_path)
            elif "msg" in file_type:
                # Action for MSG file type
                a_content = read_msg(a_file_path)
            elif "txt" in file_type:
                with open(a_file_path, "r") as file:
                    # Read the content of the file
                    a_content = file.read()
            elif "pptx" in file_type:
                a_content = read_pptx(a_file_path)
        except Exception as e:
            print("Errow with Reading " + a_file_path)
            print(e)

    return a_content


def generate_filename(output_dir):
    """
    Generate a filename output based on the current date

    @param output_dir: The directory to base the filename on.
    """
    current_date = datetime.now().strftime("%Y%m%d_%H%M%S")  # Format as %Y%m%d_%H%M%S
    filename = f"{output_dir}/selected_word_counter_output_{current_date}"
    return filename


def replace_last_slash(path, replacement=""):
    """
    This function is used for testing.

    @param path: path to be split.
    @param replacement: value to replace / with.
    """
    parts = path.rsplit("/", 1)
    return replacement.join(parts)


def delete_directory(path):
    try:
        shutil.rmtree(path)
        print(f"Directory '{path}' and all its contents were deleted successfully.")
    except FileNotFoundError:
        print(f"Directory '{path}' not found.")
    except Exception as e:
        print(f"An error occurred: {e}")
